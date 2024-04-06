import io
import tkinter as tk
from tkinter import ttk
import tkinter.messagebox
import os
import sys
from PIL import Image, ImageTk
from PyQt6 import QtWidgets
import requests
from canvasapi import Canvas
import re
from bs4 import BeautifulSoup
import hashlib
from ai_tools import generate_answer, generate_flashcards
import threading
from docx import Document
from pptx import Presentation
import fitz  # PyMuPDF
from dotenv import load_dotenv
from sniping import MyWidget
import os
load_dotenv()

API_URL = "https://q.utoronto.ca"
API_KEY = os.getenv('API_KEY')

canvas = Canvas(API_URL, API_KEY)


def select_course(courses):
    print("Please select a course:")
    for i, course in enumerate(courses):
        try:
            print(f"{i + 1}: {course.name} ({course.course_code})")
        except AttributeError:
            print(f"Course attributes not found for course index {i + 1}")

    selection = input("Enter the number of the course you want to select: ")
    try:
        selected_course = courses[int(selection) - 1]
        return selected_course
    except (IndexError, ValueError):
        print("Invalid selection. Exiting.")
        sys.exit(1)


def download_file(file, module_folder):
    try:
        response = requests.get(file.url, allow_redirects=True)
        response.raise_for_status()

        content_disp = response.headers.get('content-disposition', '')
        file_name = content_disp.split('filename=')[-1].strip('"')
        if not file_name:
            file_name = os.path.basename(file.url)

        file_path = os.path.join(module_folder, file_name)

        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        with open(file_path, 'wb') as f:
            f.write(response.content)
        print(f"Downloaded {file_name} to {file_path}")
    except requests.exceptions.RequestException as e:
        print(f"Failed to download {file.url}. Reason: {e}")


def download_page(course, item, path):
    try:
        page = course.get_page(item.page_url)
        page_title = sanitize_filename(page.title)
        page_content = clean_html(page.body)
        file_path = os.path.join(path, f"{page_title}.txt")

        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(page_content)
        print(f"Downloaded page: {page_title}")
    except Exception as e:
        print(f"Failed to download page: {item.title}. Reason: {e}")


def handle_subheader(item, path_stack):

    subheader_title = sanitize_filename(item.title)
    # Check the indentation level and adjust stack accordingly
    target_depth = item.indent if item.indent is not None else 0
    while len(path_stack) > target_depth + 1:  # +1 because the module root is the base
        path_stack.pop()  # Remove deeper or same-level directories
    # Current path is now the path of the subheader's parent
    current_path = path_stack[-1]
    subheader_dir_path = os.path.join(current_path, subheader_title)
    os.makedirs(subheader_dir_path, exist_ok=True)
    # Add the new subheader path to the stack
    path_stack.append(subheader_dir_path)
    print(f"Created directory for SubHeader: {item.title}")


def download_quiz(course, item, path):
    quiz_id = item.content_id
    try:
        quiz = course.get_quiz(quiz_id)
        quiz_title = sanitize_filename(quiz.title)
        quiz_file_path = os.path.join(path, f"{quiz_title}.txt")
        with open(quiz_file_path, 'w') as f:
            questions = quiz.get_questions()
            for question in questions:
                f.write(f"Question: {question.question_text}\n")
                for answer in question.answers:
                    f.write(f"- {answer['text']}\n\n")
        print(f"Downloaded quiz: {quiz_title}")
    except Exception as e:
        print(f"Failed to download quiz: {item.title}. Reason: {e}")


def download_assignment(course, item, path):
    assignment_id = item.content_id
    try:
        assignment = course.get_assignment(assignment_id)
        assignment_content = clean_html(assignment.description)
        valid_file_name = sanitize_filename(assignment.name)
        with open(os.path.join(path, f"{valid_file_name}.txt"), 'w') as f:
            f.write(assignment_content)
        print(f"Downloaded assignment: {assignment.name}")
    except Exception as e:
        print(f"Failed to download assignment: {item.title}. Reason: {e}")


def clean_html(html_content):
    soup = BeautifulSoup(html_content, 'html.parser')

    # Convert hyperlinks to Markdown-style links
    for a in soup.find_all('a', href=True):
        a.replace_with(f"[{a.get_text()}]({a['href']})")

    # Extract and return the cleaned text
    return soup.get_text(separator="\n")


def ensure_length(filename, max_length=255):
    # Shortens the filename, preserving the extension, if it's too long
    if len(filename) <= max_length:
        return filename
    # Split the filename into name and extension
    file_parts = os.path.splitext(filename)
    # Get a hash of the filename
    filename_hash = hashlib.sha256(filename.encode()).hexdigest()[
        :8]  # Short hash
    # Truncate the original name, add hash and preserve the file extension
    truncated_name = file_parts[0][:max_length -
                                   len(file_parts[1]) - len(filename_hash) - 1] + "_" + filename_hash
    return truncated_name + file_parts[1]

# Ensure this new function is called whenever files or directories are created


def sanitize_filename(name):
    return ensure_length(re.sub(r'[<>:"/\\|?*\x00-\x1F]', '_', name))


def save_external_url(item, path):
    external_url = item.external_url
    external_url_title = sanitize_filename(item.title)
    file_path = os.path.join(path, f"{external_url_title}.txt")
    try:
        with open(file_path, 'w') as f:
            f.write(external_url)
        print(f"Saved external URL: {external_url_title}")
    except Exception as e:
        print(
            f"Failed to save external URL '{external_url_title}'. Reason: {e}")


def download_announcements(course, course_folder):
    announcements_folder = os.path.join(course_folder, "Announcements")
    # Ensure the folder exists
    os.makedirs(announcements_folder, exist_ok=True)
    try:
        announcements = course.get_discussion_topics(only_announcements=True)
        for announcement in announcements:
            title = sanitize_filename(announcement.title)
            message = clean_html(announcement.message)
            announcement_path = os.path.join(
                announcements_folder, f"{title}.txt")
            with open(announcement_path, 'w', encoding='utf-8') as f:
                f.write(message)
            print(f"Downloaded announcement: {title}")
    except Exception as e:
        print(f"Failed to download announcements. Reason: {e}")


def print_and_download_course_details(course):
    course_name = course.name
    course_folder = os.path.join(os.getcwd(), sanitize_filename(course_name))
    os.makedirs(course_folder, exist_ok=True)

    print(f"\nCourse Name: {course_name}")
    print(f"Course Code: {course.course_code}\n")
    download_announcements(course, course_folder)

    modules = course.get_modules()
    if not list(modules):
        print("No modules found.")
        # Get the root files folder
        root_folder = None
        for folder in course.get_folders():
            if folder.parent_folder_id is None:
                root_folder = folder
                break

        if root_folder:
            # Use the Canvas API to explore folder hierarchy

            def get_files_in_folder_recursive(folder, current_path):
                folder_files_url = f"{API_URL}/api/v1/folders/{folder.id}/files"
                try:
                    file_response = requests.get(folder_files_url, headers={
                                                 'Authorization': 'Bearer {}'.format(API_KEY)})
                    # This will raise an HTTPError if the HTTP request returned an unsuccessful status code
                    file_response.raise_for_status()
                    files_json = file_response.json()

                    for file_data in files_json:
                        file_obj = canvas.get_file(file_data['id'])
                        download_file(file_obj, current_path)

                    subfolders = folder.get_folders()
                    for subfolder in subfolders:
                        subfolder_name = sanitize_filename(subfolder.name)
                        new_path = os.path.join(current_path, subfolder_name)
                        os.makedirs(new_path, exist_ok=True)
                        get_files_in_folder_recursive(subfolder, new_path)

                except requests.exceptions.HTTPError as e:
                    # Here we catch HTTP errors, which include the 403 Forbidden
                    print(
                        f"HTTP Request failed: {e.response.status_code} {e.response.reason} for url: {folder_files_url}")
                    if e.response.status_code == 403:
                        print(
                            "It looks like we don't have access to this resource. Skipping...")
                except requests.exceptions.RequestException as e:
                    # This catches any other exceptions that requests might raise
                    print(
                        f"Failed to get files in folder: {folder.name}. Reason: {e}")
                except Exception as e:
                    # This will catch any other exceptions
                    print(f"An unexpected error occurred: {e}")

            # Start recursive file and folder download process
            root_folder_name = sanitize_filename(root_folder.name)
            root_folder_path = os.path.join(course_folder, root_folder_name)
            os.makedirs(root_folder_path, exist_ok=True)
            try:
                get_files_in_folder_recursive(root_folder, root_folder_path)
            except Exception as e:
                print(f"Failed to download files. Reason: {e}")

        else:
            print("Root folder not found, downloading files from course root.")
            for file in course.get_files():
                download_file(file, course_folder)
    else:
        print(f"Downloading modules for course: {course_name}")
        for module in modules:
            module_name = module.name
            print(f"\nDownloading module: {module_name}")

            module_folder = os.path.join(
                course_folder, sanitize_filename(module_name))
            os.makedirs(module_folder, exist_ok=True)
            # Initialize the stack with the module folder
            path_stack = [module_folder]

            for module_item in module.get_module_items():
                if module_item.type == "File":
                    file = canvas.get_file(module_item.content_id)
                    download_file(file, path_stack[-1])
                elif module_item.type == "Page":
                    download_page(course, module_item, path_stack[-1])
                elif module_item.type == "Quiz":
                    download_quiz(course, module_item, path_stack[-1])
                elif module_item.type == "Assignment":
                    download_assignment(course, module_item, path_stack[-1])
                elif module_item.type == "ExternalUrl":
                    save_external_url(module_item, path_stack[-1])
                elif module_item.type == "SubHeader":

                    handle_subheader(module_item, path_stack)
                else:
                    print(f"Item type {module_item.type} not handled.")

            # Reset the path_stack for the next module
            path_stack = [module_folder]
    return course_folder


def flashcardGUI(flashcards: dict):
    class FlashcardsApp(tk.Tk):
        def __init__(self, flashcard_data):
            super().__init__()

            self.flashcards = flashcard_data["flashcards"]
            self.current_card = 0

            self.title("Flashcards")
            self.geometry("800x400")

            self.question_label = tk.Label(self, text=self.flashcards[self.current_card]["question"],
                                           font=("Arial", 20, "bold"), wraplength=700)
            self.question_label.pack(pady=50)

            self.show_answer_button = tk.Button(self, text="Show Answer",
                                                command=self.show_answer, font=("Arial", 15))
            self.show_answer_button.pack()

            self.next_button = tk.Button(self, text="Next", command=self.next_card,
                                         font=("Arial", 15), state="disabled")
            self.next_button.pack(pady=20)

        def show_answer(self):
            self.show_answer_button.config(state="disabled")

            self.answer_label = tk.Label(self, text=self.flashcards[self.current_card]["answer"],
                                         font=("Arial", 15), wraplength=700)
            self.answer_label.pack(pady=10)

            if self.current_card < len(self.flashcards) - 1:
                self.next_button.config(state="normal")

        def next_card(self):
            self.current_card += 1
            self.question_label.config(
                text=self.flashcards[self.current_card]["question"])
            self.answer_label.destroy()

            self.show_answer_button.config(state="normal")
            self.next_button.config(state="disabled")

    app = FlashcardsApp(flashcards)
    app.mainloop()


class CourseApp(tk.Tk):
    def __init__(self, course_folder, context=""):
        super().__init__()
        self.context = context
        self.course_folder = course_folder
        self.open_tabs = {}
        self.image = "snips/snip.png"  # Attribute to store the image path

        self.title("Quercus Assistant")
        self.geometry("1920x1080")

        self.create_widgets()
        self.update_treeview(self.course_folder)

    def create_widgets(self):
        self.notebook = ttk.Notebook(self)
        self.create_main_frame()
        self.notebook.pack(expand=True, fill="both")

    def create_main_frame(self):
        main_frame = ttk.Frame(self.notebook)
        main_frame.columnconfigure(0, weight=1)
        main_frame.columnconfigure(1, weight=1)
        main_frame.columnconfigure(2, weight=1)
        main_frame.rowconfigure(0, weight=1)

        # File selection column
        self.treeview = ttk.Treeview(main_frame, selectmode="extended")
        self.treeview_scroll = ttk.Scrollbar(
            main_frame, orient="vertical", command=self.treeview.yview
        )
        self.treeview.configure(yscrollcommand=self.treeview_scroll.set)
        self.treeview.grid(row=0, column=0, sticky="nsew")
        self.treeview_scroll.grid(row=0, column=1, sticky="ns")

        # Chat interface column
        self.chat_text = tk.Text(main_frame, wrap="word")
        self.chat_text_scroll = ttk.Scrollbar(
            main_frame, orient="vertical", command=self.chat_text.yview
        )
        self.chat_text.configure(yscrollcommand=self.chat_text_scroll.set)
        self.chat_text.grid(row=0, column=2, padx=(10, 0), sticky="nsew")
        self.chat_text_scroll.grid(row=0, column=3, sticky="ns")

        self.chat_input = ttk.Entry(main_frame)
        self.chat_input.grid(row=1, column=2, padx=(10, 0), sticky="ew")
        self.chat_input.bind("<Return>", self.chat_send_event)
        # Bind double-click event
        self.treeview.bind("<Double-1>", self.on_file_double_click)

        self.chat_send_button = ttk.Button(
            main_frame, text="Send", command=self.send_chat)
        self.chat_send_button.grid(row=1, column=3, padx=(0, 10), sticky="ew")

        self.flashcards = ttk.Button(
            main_frame, text="generate flashcards", command=self.createFlashcards)
        self.flashcards.grid(row=1, column=4, padx=(0, 10), sticky="ew")

        self.chat_image_label = tk.Label(main_frame)
        self.chat_image_label.grid(row=2, column=2, padx=(10, 0), sticky="ew")

        self.notebook.add(main_frame, text="Main")
        self.snip_button = ttk.Button(self, text="Snip", command=self.snip)
        self.snip_button.place(relx=1.0, rely=0.0, anchor='ne')

    def snip(self):
        app = QtWidgets.QApplication(sys.argv)
        window = MyWidget()
        window.show()
        app.aboutToQuit.connect(app.deleteLater)
        app.exec()
        self.update_chat_image("snips/snip.png")

    def createFlashcards(self):
        if not self.treeview.selection():
            tkinter.messagebox.showinfo(
                "Info", "Please select files from the treeview.")
            return
        complete_paths = []
        for file_item in self.treeview.selection():
            complete_paths.append(os.path.join(
                self.course_folder, self.get_filepath(file_item, self.course_folder)))
        tkinter.messagebox.showinfo(
            "Info", "Generating flashcards. Please wait...")
        flashcard_dict = generate_flashcards(
            complete_paths, self.context)

        # Call the flashcardGUI function with the flashcard_dict
        flashcardGUI(flashcard_dict)

        # self.send_response("bot", bot_response)

    def send_response(self, response_type, message):
        color_map = {
            "user": "green",
            "bot": "white",
            "system": "purple"
        }
        font_map = {
            "user": ("Arial", 20),
            "bot": ("Arial", 20, "bold"),
            "system": ("Arial", 20, "bold")
        }

        self.chat_text.tag_configure(
            response_type, foreground=color_map[response_type], font=font_map[response_type])
        self.chat_text.insert(
            "end", f"{response_type.upper()}: {message}\n", response_type)

    def get_filepath(self, tree_item, base_folder):
        parent = self.treeview.parent(tree_item)
        if parent:
            return os.path.join(self.get_filepath(parent, base_folder), self.treeview.item(tree_item, option="text"))
        else:
            return self.treeview.item(tree_item, option="text")

    def chat_send_event(self, event):
        self.send_chat()

    def send_chat(self):
        chat_message = self.chat_input.get()

        if not chat_message and not self.image:
            return

        display_message = chat_message
        if self.image:
            display_message += f"\n[Image: {self.image}]"

        self.send_response("user", display_message)
        self.chat_input.delete(0, "end")

        # Clear the image after sending the message
        self.chat_image_label.config(image='')
        self.clear_chat_image()

        self.chat_text.see("end")
        self.update_idletasks()

        threading.Thread(target=self.generate_response,
                         args=(chat_message, "snips/snip.png")).start()

    def clear_chat_image(self):
        """ Clear the chat image """
        self.chat_image_label.config(image='')

    def update_chat_image(self, image_path):
        """ Update the chat with an image based on the provided path """
        if os.path.isfile(image_path):
            self.image = ImageTk.PhotoImage(Image.open(image_path))
            self.chat_image_label.config(image=self.image)
            self.chat_image_label.image = self.image  # Keep a reference
        else:
            print("Image file does not exist.")

    def generate_response(self, chat_message, image=None):
        # Get complete paths of selected files
        complete_paths = []
        for file_item in self.treeview.selection():
            complete_paths.append(os.path.join(
                self.course_folder, self.get_filepath(file_item, self.course_folder)))

        bot_response = generate_answer(
            chat_message, complete_paths, self.context, image=image)
        self.send_response("bot", bot_response)

    def insert_files_recursively(self, folder, root_node):
        for item in os.listdir(folder):
            item_path = os.path.join(folder, item)
            if os.path.isdir(item_path):
                directory_node = self.treeview.insert(
                    root_node, "end", text=item, open=True,
                )
                self.insert_files_recursively(item_path, directory_node)
            else:
                self.treeview.insert(root_node, "end", text=item)

    def update_treeview(self, folder):
        self.course_folder = folder
        self.treeview.delete(*self.treeview.get_children())
        self.insert_files_recursively(self.course_folder, "")

    def on_file_double_click(self, event):
        selected_item = self.treeview.selection()[0]
        file_path = os.path.join(self.course_folder, self.get_filepath(
            selected_item, self.course_folder))
        if os.path.isfile(file_path):
            self.open_or_focus_file_tab(file_path)

    def open_or_focus_file_tab(self, file_path):
        tab_name = os.path.basename(file_path)
        if tab_name in self.open_tabs:
            self.notebook.select(self.open_tabs[tab_name])
        else:
            self.open_file_in_tab(file_path, tab_name)

    def open_file_in_tab(self, file_path, tab_name):
        new_tab = ttk.Frame(self.notebook)
        self.notebook.add(new_tab, text=tab_name)
        self.notebook.tab(new_tab, compound=tk.TOP)
        self.open_tabs[tab_name] = new_tab

        close_button = ttk.Button(
            new_tab, text="✕", command=lambda: self.close_tab(tab_name))
        close_button.pack(side="right", anchor="ne")

        if file_path.endswith(".pdf"):
            self.display_pdf(file_path, new_tab)
        elif file_path.endswith(".docx"):
            self.display_docx(file_path, new_tab)
        elif file_path.endswith(".pptx"):
            self.display_pptx(file_path, new_tab)
        else:
            self.display_text_file(file_path, new_tab)

    def close_tab(self, tab_name):
        if tab_name in self.open_tabs:
            tab = self.open_tabs.pop(tab_name)
            self.notebook.forget(tab)
            tab.destroy()

    def display_text_file(self, file_path, tab):
        text_widget = tk.Text(tab, wrap="word")
        text_widget.pack(expand=True, fill="both")
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                text_widget.insert("1.0", file.read())
        except Exception as e:
            text_widget.insert("1.0", f"Error opening file: {e}")

    def display_pdf(self, file_path, tab):
        doc = fitz.open(file_path)
        pdf_viewer = ttk.Notebook(tab)
        pdf_viewer.pack(expand=True, fill="both")

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            img_data = pix.tobytes("ppm")
            img = Image.open(io.BytesIO(img_data))
            photo = ImageTk.PhotoImage(image=img)

            canvas = tk.Canvas(pdf_viewer, width=pix.width, height=pix.height)
            canvas.create_image(0, 0, image=photo, anchor="nw")
            canvas.image = photo  # Keep a reference.
            pdf_viewer.add(canvas, text=f"Page {page_num + 1}")

    def display_docx(self, file_path, tab):
        document = Document(file_path)
        text_widget = tk.Text(tab, wrap="word")
        text_widget.pack(expand=True, fill="both")
        for para in document.paragraphs:
            text_widget.insert("end", para.text + "\n")

    def display_pptx(self, file_path, tab):
        presentation = Presentation(file_path)
        text_widget = tk.Text(tab, wrap="word")
        text_widget.pack(expand=True, fill="both")
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_widget.insert("end", shape.text + "\n")


def main():
    print(f"Current working directory: {os.getcwd()}")
    try:
        courses = canvas.get_courses()
        if not courses:
            print("No courses found. Exiting.")
            sys.exit(0)

        selected_course = select_course(courses)
        # print_and_download_course_details(selected_course)
        app = CourseApp(print_and_download_course_details(selected_course))
        app.mainloop()
        # app = CourseApp("CCT109H5 F LEC0101 & LEC0102",
        #                 context=f" The student is taking CCT109H5 F LEC0101 & LEC0102, offered ut the University of Toronto Mississauga.")
        app = CourseApp(print_and_download_course_details(selected_course),
                        context=f" The student is taking {selected_course}, offered ut the University of Toronto Mississauga.")
        app.mainloop()
    except Exception as e:
        print(f"An unexpected error occurred: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
